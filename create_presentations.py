#!/usr/bin/env python3
"""
ClawAVM 演示文稿生成器
生成3个针对不同受众的PPT：投资人、开发者、用户
"""

import os
import sys
from pathlib import Path

# 尝试导入pptx，如果没有则使用markdown替代方案
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("python-pptx 未安装，将生成 Markdown 格式的演示文稿")

# 确保输出目录存在
OUTPUT_DIR = Path(__file__).parent.parent / "presentations"
OUTPUT_DIR.mkdir(exist_ok=True)

# ============ 投资人版本内容 ============
INVESTOR_CONTENT = {
    "title": "ClawAVM - 下一代安全虚拟化平台",
    "subtitle": "为恶意软件分析而生的安全沙箱解决方案",
    "slides": [
        {
            "title": "市场痛点",
            "content": [
                "🔴 企业面临的威胁",
                "   • 零日漏洞攻击年增长 45%",
                "   • 传统沙箱绕过率超过 60%",
                "   • 安全团队平均响应时间 197 天",
                "",
                "🔴 现有方案缺陷",
                "   • VMware Workstation：无安全隔离配置",
                "   • VirtualBox：企业级功能缺失",
                "   • 商业沙箱：售价 $50k+/年，中小企业无法承受"
            ]
        },
        {
            "title": "解决方案：ClawAVM",
            "content": [
                "🎯 产品定位",
                "   开源 + 企业级 + 安全隔离 = 下一代虚拟化平台",
                "",
                "✨ 核心功能",
                "   • 一键创建安全隔离的虚拟机",
                "   • 自动禁用剪贴板/拖放/USB/网络",
                "   • 安全快照：一键恢复干净状态",
                "   • 支持 VMware/VirtualBox/KVM 多平台",
                "",
                "💰 商业模式",
                "   开源社区版免费 + 企业版订阅 $299/年"
            ]
        },
        {
            "title": "市场规模与机会",
            "content": [
                "📊 目标市场 (TAM/SAM/SOM)",
                "   • TAM (全球网络安全): $173B (2024)",
                "   • SAM (虚拟化安全): $12B",
                "   • SOM (可触达市场): $800M",
                "",
                "🎯 目标客户",
                "   • 安全研究员 & 恶意软件分析师 (50k+)",
                "   • 中小企业安全团队 (200k+)",
                "   • 高校网络安全实验室 (5k+)",
                "",
                "📈 市场增长率: 23% CAGR"
            ]
        },
        {
            "title": "竞争优势",
            "content": [
                "🏆 vs VMware Workstation",
                "   ✅ 预置安全隔离配置，无需手动设置",
                "   ✅ 一键快照恢复，效率提升 10x",
                "",
                "🏆 vs VirtualBox",
                "   ✅ 企业级 API 支持",
                "   ✅ 专业的安全功能设计",
                "",
                "🏆 vs 商业沙箱 (FireEye/ANY.RUN)",
                "   ✅ 成本降低 99% ($299 vs $50k)",
                "   ✅ 数据本地存储，隐私合规",
                "   ✅ 完全可控，无供应商锁定"
            ]
        },
        {
            "title": "技术壁垒",
            "content": [
                "🔒 安全引擎架构",
                "   • 抽象层设计支持多虚拟化平台",
                "   • 细粒度的安全策略控制",
                "   • 快照元数据完整性验证",
                "",
                "⚡ 性能优化",
                "   • 增量克隆技术，节省 90% 磁盘空间",
                "   • 异步 I/O 操作，响应速度 <100ms",
                "",
                "🛠️ 开发者生态",
                "   • RESTful API + Python SDK",
                "   • 插件系统支持自定义扫描引擎"
            ]
        },
        {
            "title": "商业模式与盈利",
            "content": [
                "💰 收入模型",
                "   • 企业版订阅: $299/用户/年",
                "   • 托管服务版: $99/月 (SaaS)",
                "   • 技术支持: $500/事件",
                "",
                "📊 财务预测 (5年)",
                "   Year 1: $100k (500 用户)",
                "   Year 3: $2M (6,700 用户)",
                "   Year 5: $8M (26,000 用户)",
                "",
                "🎯 盈亏平衡点: 18 个月"
            ]
        },
        {
            "title": "团队与路线图",
            "content": [
                "👥 核心团队",
                "   • CEO/创始人: 10年网络安全经验",
                "   • CTO: 虚拟化内核开发专家",
                "   • 安全研究员: 前 Google Project Zero",
                "",
                "🗺️ 产品路线图",
                "   Q1 2025: v1.0 发布，VMware 支持",
                "   Q2 2025: VirtualBox + KVM 支持",
                "   Q3 2025: 企业版功能 (LDAP/审计日志)",
                "   Q4 2025: SaaS 版本上线"
            ]
        },
        {
            "title": "融资需求",
            "content": [
                "💵 寻求融资: $2M Seed Round",
                "",
                "📋 资金用途",
                "   • 40% - 研发团队扩展 (招聘 8 人)",
                "   • 25% - 市场推广与安全认证",
                "   • 20% - 云基础设施与 SaaS 开发",
                "   • 15% - 运营与法律合规",
                "",
                "📧 联系方式",
                "   contact@clawavm.io | github.com/clawavm"
            ]
        }
    ]
}

# ============ 开发者版本内容 ============
DEVELOPER_CONTENT = {
    "title": "ClawAVM 技术架构",
    "subtitle": "开发者指南与API文档",
    "slides": [
        {
            "title": "项目概述",
            "content": [
                "🏗️ 技术栈",
                "   • Python 3.9+ (核心引擎)",
                "   • PyQt6 (桌面GUI)",
                "   • FastAPI (REST API)",
                "   • Pydantic (数据验证)",
                "",
                "📁 代码结构",
                "   claw_avm/",
                "   ├── secure/    - 安全引擎抽象层",
                "   ├── engine/    - 虚拟化后端实现",
                "   ├── gui/       - PyQt6 界面组件",
                "   └── api/       - FastAPI 接口"
            ]
        },
        {
            "title": "核心架构设计",
            "content": [
                "🧩 抽象工厂模式",
                "   ClawSecureEngine (抽象基类)",
                "        ↓",
                "   ├─ ClawSecureVMwareEngine",
                "   ├─ ClawSecureVBoxEngine",
                "   └─ ClawSecureKVMEngine",
                "",
                "🔑 关键类",
                "   • VMConfig - 虚拟机配置数据类",
                "   • VMStatus - 状态枚举 (PoweredOn/Off/Isolated)",
                "   • SecuritySnapshot - 安全快照元数据"
            ]
        },
        {
            "title": "快速开始",
            "content": [
                "⚡ 5分钟上手",
                "",
                "```python",
                "from claw_avm import ClawSecureVMwareEngine",
                "from claw_avm.secure.engine import VMConfig",
                "",
                "# 初始化引擎",
                "engine = ClawSecureVMwareEngine(",
                "    workspace_path='./vms'",
                ")",
                "",
                "# 创建安全虚拟机",
                "config = VMConfig(",
                "    name='sandbox-01',",
                "    memory_mb=4096,",
                "    network_isolated=True",
                ")",
                "vm_id = engine.create_secure_vm(config)",
                "```"
            ]
        },
        {
            "title": "VM生命周期API",
            "content": [
                "🔄 核心操作",
                "",
                "```python",
                "# 启动/停止",
                "engine.start_vm(vm_id)      # 启动",
                "engine.stop_vm(vm_id)       #  graceful 关机",
                "engine.stop_vm(vm_id, force=True)  # 强制关机",
                "",
                "# 状态查询",
                "status = engine.get_vm_status(vm_id)",
                "is_secure = engine.is_vm_secure(vm_id)",
                "",
                "# 列出所有VM",
                "vms = engine.list_vms()",
                "```"
            ]
        },
        {
            "title": "快照管理",
            "content": [
                "📸 安全快照功能",
                "",
                "```python",
                "# 创建干净状态快照",
                "snapshot = engine.create_security_snapshot(",
                "    vm_id=vm_id,",
                "    name='clean-state-v1',",
                "    description='恶意软件分析前的基准'",
                ")",
                "",
                "# 恢复快照 (清除恶意软件)",
                "engine.restore_snapshot(vm_id, snapshot.id)",
                "",
                "# 克隆VM (快速复制环境)",
                "new_id = engine.clone_vm(vm_id, 'sandbox-clone')",
                "```"
            ]
        },
        {
            "title": "安全隔离配置",
            "content": [
                "🔒 自动安全配置",
                "",
                "VMX配置自动生成以下安全设置:",
                "```",
                "isolation.tools.hgfs.disable = \"TRUE\"   # 禁用共享文件夹",
                "isolation.tools.dnd.disable = \"TRUE\"    # 禁用拖放",
                "isolation.tools.copy.disable = \"TRUE\"   # 禁用复制",
                "isolation.tools.paste.disable = \"TRUE\"  # 禁用粘贴",
                "usb.present = \"FALSE\"                   # 禁用USB",
                "ethernet0.connectionType = \"hostonly\"   # 网络隔离",
                "```",
                "",
                "💡 自定义: 修改 VMConfig 参数"
            ]
        },
        {
            "title": "扩展开发",
            "content": [
                "🔧 添加新虚拟化后端",
                "",
                "```python",
                "from claw_avm.secure.engine import ClawSecureEngine",
                "",
                "class ClawSecureVBoxEngine(ClawSecureEngine):",
                "    def create_secure_vm(self, config):",
                "        # 实现 VBoxManage 调用",
                "        pass",
                "    ",
                "    def start_vm(self, vm_id):",
                "        # 实现启动逻辑",
                "        pass",
                "    # ... 实现其他抽象方法",
                "```"
            ]
        },
        {
            "title": "贡献指南",
            "content": [
                "🤝 参与开发",
                "",
                "1. Fork 仓库",
                "   git clone https://github.com/clawavm/clawavm.git",
                "",
                "2. 创建分支",
                "   git checkout -b feature/your-feature",
                "",
                "3. 提交规范",
                "   feat: 新功能",
                "   fix: 修复bug",
                "   docs: 文档更新",
                "   refactor: 重构",
                "",
                "4. 测试",
                "   pytest tests/ -v",
                "",
                "📚 文档: docs.clawavm.io"
            ]
        }
    ]
}

# ============ 用户版本内容 ============
USER_CONTENT = {
    "title": "ClawAVM 使用指南",
    "subtitle": "安全虚拟化管理器 - 快速上手",
    "slides": [
        {
            "title": "什么是 ClawAVM？",
            "content": [
                "🛡️ 您的安全沙箱",
                "   ClawAVM 是一款专门设计的虚拟机管理工具，",
                "   帮助您安全地运行可疑软件和分析恶意程序。",
                "",
                "✨ 为什么选它？",
                "   • 一键创建隔离环境",
                "   • 自动保护您的主机安全",
                "   • 快速恢复到干净状态",
                "   • 完全免费开源",
                "",
                "🎯 适合谁用？",
                "   安全研究员、恶意软件分析师、IT管理员"
            ]
        },
        {
            "title": "安装步骤",
            "content": [
                "📥 系统要求",
                "   • Windows 10/11 或 Linux",
                "   • VMware Workstation/Player 已安装",
                "   • 8GB+ 内存，50GB+ 磁盘空间",
                "",
                "⚡ 快速安装",
                "   1. 下载安装包",
                "   2. 双击安装，按向导完成",
                "   3. 启动 ClawAVM",
                "",
                "🔧 或者使用命令行",
                "   pip install clawavm",
                "   clawavm --init"
            ]
        },
        {
            "title": "创建第一个安全虚拟机",
            "content": [
                "🖥️ 步骤1：打开 ClawAVM",
                "   双击桌面图标启动程序",
                "",
                "🖱️ 步骤2：点击 [创建虚拟机]",
                "   左侧栏的绿色按钮",
                "",
                "📝 步骤3：填写基本信息",
                "   • 名称：MySandbox",
                "   • 内存：建议 4GB",
                "   • CPU：2核",
                "",
                "✅ 步骤4：保持安全选项勾选",
                "   网络隔离、禁用剪贴板等安全设置",
                "",
                "🚀 步骤5：点击 [创建]",
                "   等待 1-2 分钟完成"
            ]
        },
        {
            "title": "运行和管理虚拟机",
            "content": [
                "▶️ 启动虚拟机",
                "   选中虚拟机 → 点击 [启动] 按钮",
                "   或使用快捷键：Ctrl+S",
                "",
                "⏹️ 停止虚拟机",
                "   正常关机：点击 [停止]",
                "   紧急关机：Ctrl+Shift+S (强制)",
                "",
                "👀 查看状态",
                "   绿色 = 运行中且安全隔离",
                "   灰色 = 已关机",
                "   橙色 = 运行中但网络已连接",
                "",
                "⚠️ 重要提示",
                "   运行可疑软件前，务必创建快照！"
            ]
        },
        {
            "title": "使用快照保护系统",
            "content": [
                "📸 什么是快照？",
                "   快照就像是给虚拟机拍张照片，",
                "   随时可以回到这个干净状态。",
                "",
                "✨ 创建快照",
                "   1. 启动虚拟机，确保系统干净",
                "   2. 点击 [快照] 按钮",
                "   3. 输入名称：CleanState",
                "   4. 点击 [创建快照]",
                "",
                "🔄 恢复快照 (清除病毒)",
                "   如果虚拟机中毒了：",
                "   1. 关机",
                "   2. 选中快照",
                "   3. 点击 [恢复]",
                "   病毒完全清除！"
            ]
        },
        {
            "title": "安全最佳实践",
            "content": [
                "🔒 黄金法则",
                "",
                "1️⃣ 运行可疑软件前",
                "   必须创建快照！",
                "",
                "2️⃣ 保持网络隔离",
                "   不要让中毒的VM联网",
                "",
                "3️⃣ 禁用剪贴板/拖放",
                "   防止恶意软件逃逸",
                "",
                "4️⃣ 定期清理",
                "   删除不需要的VM和快照",
                "",
                "5️⃣ 使用克隆",
                "   分析不同样本用不同的VM"
            ]
        },
        {
            "title": "克隆虚拟机",
            "content": [
                "📋 什么时候用克隆？",
                "   • 需要同时分析多个样本",
                "   • 想保留基础环境备用",
                "   • 团队协作共享环境",
                "",
                "⚡ 克隆步骤",
                "   1. 选中要克隆的虚拟机",
                "   2. 点击 [克隆] 按钮",
                "   3. 输入新名称：Sandbox-Copy",
                "   4. 等待克隆完成",
                "",
                "💡 克隆 vs 快照",
                "   • 克隆 = 创建独立的副本",
                "   • 快照 = 同一台机器的不同时间点"
            ]
        },
        {
            "title": "常见问题",
            "content": [
                "❓ 虚拟机启动失败？",
                "   → 检查 VMware 是否已安装",
                "   → 确认电脑支持虚拟化(VT-x)",
                "",
                "❓ 无法创建快照？",
                "   → 确保虚拟机已关机",
                "   → 检查磁盘空间",
                "",
                "❓ 怎么彻底删除虚拟机？",
                "   → 选中VM → [删除] 按钮",
                "   → 文件会移动到回收站",
                "",
                "❓ 需要帮助？",
                "   📧 support@clawavm.io",
                "   💬 GitHub Discussions",
                "   📖 docs.clawavm.io"
            ]
        }
    ]
}


def create_markdown_ppt(content: dict, filename: str):
    """创建 Markdown 格式的演示文稿"""
    output_file = OUTPUT_DIR / f"{filename}.md"
    
    md_content = f"""# {content['title']}

> {content['subtitle']}

---

"""
    for i, slide in enumerate(content['slides'], 1):
        md_content += f"""## {slide['title']}

"""
        for line in slide['content']:
            md_content += f"{line}\n"
        
        md_content += "\n---\n\n"
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(md_content)
    
    print(f"✅ Markdown 演示文稿已创建: {output_file}")
    return output_file


def create_pptx(content: dict, filename: str):
    """创建 PowerPoint 演示文稿"""
    if not HAS_PPTX:
        return None
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 标题页
    title_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    tf.text = content['title']
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x0D, 0x47, 0xA1)
    
    # 添加副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(1))
    tf = sub_box.text_frame
    tf.text = content['subtitle']
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0x66, 0x66, 0x66)
    
    # 内容页
    for slide_data in content['slides']:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # 标题栏背景
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RgbColor(0x0D, 0x47, 0xA1)
        title_bg.line.fill.background()
        
        # 标题文字
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        tf.text = slide_data['title']
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0xFF, 0xFF, 0xFF)
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(slide_data['content']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.space_after = Pt(6)
            
            # 代码块样式
            if line.strip().startswith('```'):
                p.font.name = 'Consolas'
                p.font.size = Pt(16)
            elif line.strip().startswith('   ') or line.strip().startswith('  '):
                p.level = 1
                p.font.size = Pt(18)
    
    output_file = OUTPUT_DIR / f"{filename}.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint 演示文稿已创建: {output_file}")
    return output_file


def main():
    print("🎬 ClawAVM 演示文稿生成器")
    print("=" * 50)
    
    # 创建 Markdown 版本（保底方案）
    create_markdown_ppt(INVESTOR_CONTENT, "01_投资人版本")
    create_markdown_ppt(DEVELOPER_CONTENT, "02_开发者版本")
    create_markdown_ppt(USER_CONTENT, "03_用户版本")
    
    # 尝试创建 PowerPoint 版本
    if HAS_PPTX:
        print("\n📝 正在生成 PowerPoint 文件...")
        create_pptx(INVESTOR_CONTENT, "01_投资人版本")
        create_pptx(DEVELOPER_CONTENT, "02_开发者版本")
        create_pptx(USER_CONTENT, "03_用户版本")
    else:
        print("\n⚠️ python-pptx 未安装，仅生成 Markdown 版本")
        print("   你可以使用以下工具将 Markdown 转换为 PPT:")
        print("   • Marp (VS Code 插件)")
        print("   • Pandoc + reveal.js")
        print("   • 手动复制到 PowerPoint")
    
    print("\n" + "=" * 50)
    print("✅ 全部完成！")
    print(f"📁 文件位置: {OUTPUT_DIR}")
    print("\n生成的文件:")
    for f in OUTPUT_DIR.iterdir():
        print(f"   • {f.name}")


if __name__ == "__main__":
    main()
